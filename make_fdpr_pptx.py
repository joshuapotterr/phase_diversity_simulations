from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── colour palette ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)   # slide background
BLUE    = RGBColor(0x1E, 0x6F, 0xC8)   # accent / headings
CYAN    = RGBColor(0x2E, 0xCC, 0xE0)   # highlight
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xD0, 0xD8, 0xE8)   # body text
ORANGE  = RGBColor(0xF0, 0x8C, 0x00)   # warning / callout
GREEN   = RGBColor(0x2E, 0xC4, 0x7A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank

# ── low-level helpers ────────────────────────────────────────────────────────

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height,
        fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shp = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.line.width = line_width
    if fill_color:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = line_width if line_width else Pt(1)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def accent_bar(slide, top=0.55, height=0.06, color=BLUE):
    box(slide, 0, top, 13.33, height, fill_color=color)

def slide_header(slide, title, subtitle=None):
    accent_bar(slide, top=0.0, height=0.55, color=NAVY)
    accent_bar(slide, top=0.55, height=0.06, color=BLUE)
    txt(slide, title, 0.35, 0.06, 12.5, 0.5,
        size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.62, 12.5, 0.35,
            size=14, color=CYAN, align=PP_ALIGN.LEFT)

def bullet_block(slide, items, left, top, width, height,
                 title=None, title_color=CYAN, bullet_color=LGRAY,
                 size=14, title_size=16, fill=None, border=None):
    if fill or border:
        box(slide, left, top, width, height,
            fill_color=fill, line_color=border or BLUE, line_width=Pt(1))
    if title:
        txt(slide, title, left+0.1, top+0.08, width-0.2, 0.3,
            size=title_size, bold=True, color=title_color)
        top += 0.38
        height -= 0.38
    txb = slide.shapes.add_textbox(
        Inches(left+0.1), Inches(top), Inches(width-0.2), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = ("• " if not item.startswith("  ") else "") + item
        run.font.size  = Pt(size)
        run.font.color.rgb = bullet_color

def code_block(slide, code_lines, left, top, width, height):
    box(slide, left, top, width, height,
        fill_color=RGBColor(0x0A,0x12,0x20),
        line_color=BLUE, line_width=Pt(1))
    txb = slide.shapes.add_textbox(
        Inches(left+0.12), Inches(top+0.1),
        Inches(width-0.24), Inches(height-0.2))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.color.rgb = CYAN
        run.font.name = "Courier New"

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
box(s, 0, 0, 13.33, 7.5, fill_color=NAVY)
box(s, 0, 2.9, 13.33, 0.08, fill_color=BLUE)
box(s, 0, 3.08, 13.33, 0.04, fill_color=CYAN)

txt(s, "FDPR Pipeline Comparison", 1.0, 1.3, 11.3, 1.2,
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Focus-Diverse Phase Retrieval: Defocus Models, Solvers & Readout Methods",
    1.0, 2.55, 11.3, 0.55, size=18, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, "fdpr_vs_AS_save.py  ·  FDPRno_phase_error.py",
    1.0, 3.25, 11.3, 0.5, size=14, color=LGRAY,
    align=PP_ALIGN.CENTER, italic=True)
txt(s, "SEAL Lab  ·  2025", 1.0, 6.8, 11.3, 0.4,
    size=12, color=LGRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is FDPR?
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "What is Focus-Diverse Phase Retrieval?",
             "Recovering wavefront aberrations from intensity images alone")

# left column
bullet_block(s, [
    "The pupil plane and focal plane are Fourier pairs",
    "A PSF encodes phase information — but cameras only measure intensity (|E|²)",
    "Phase is lost: we must infer it from multiple intensity images",
    "",
    "Focus diversity = deliberately defocus the camera to 'diverse' positions",
    "Each defocused image adds a known phase offset → more constraints",
    "Iterate until a phase is consistent with ALL measured intensities",
], left=0.3, top=0.75, width=6.5, height=5.5,
   title="Core Idea", fill=RGBColor(0x10,0x20,0x38), border=BLUE)

# right column — GS loop
bullet_block(s, [
    "1.  Start with a random phase estimate",
    "2.  Build complex focal field:  E = A_focus · e^(iφ)",
    "3.  FFT → pupil plane",
    "4.  Add known defocus → defocused pupil",
    "5.  FFT → defocused focal plane",
    "6.  FORCE amplitude to match measurement, keep phase",
    "7.  Invert: remove defocus, go back to focus",
    "8.  FORCE focused amplitude constraint",
    "9.  Repeat for each defocused image",
    "10. Iterate until convergence (MSE plateaus)",
], left=7.0, top=0.75, width=6.0, height=5.5,
   title="Gerchberg-Saxton Loop", fill=RGBColor(0x10,0x20,0x38), border=CYAN)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Two Scripts Overview
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "Two Scripts, Two Goals")

# left card
box(s, 0.3, 0.75, 5.9, 5.8, fill_color=RGBColor(0x08,0x18,0x30), line_color=CYAN, line_width=Pt(2))
txt(s, "FDPRno_phase_error.py", 0.4, 0.82, 5.7, 0.4,
    size=16, bold=True, color=CYAN)
txt(s, "Baseline / Noise Floor", 0.4, 1.22, 5.7, 0.3,
    size=12, color=LGRAY, italic=True)
bullet_block(s, [
    "Single fixed pipeline",
    "NO injected aberration (flat phase)",
    "6 defocus positions (±4, ±8, ±12 mm)",
    "Goal: measure algorithm's intrinsic noise floor",
    "",
    "Solver:   FocusDiversePhaseRetrieval",
    "Defocus:  Zernike pupil-plane",
    "Readout:  mft_rev → resize → mask",
    "",
    "Single run — no Monte Carlo sweep",
    "Truth = zero phase; any residual = FDPR error",
], left=0.4, top=1.55, width=5.7, height=4.8, size=13, bullet_color=LGRAY)

# right card
box(s, 6.8, 0.75, 6.2, 5.8, fill_color=RGBColor(0x08,0x18,0x30), line_color=ORANGE, line_width=Pt(2))
txt(s, "fdpr_vs_AS_save.py", 6.9, 0.82, 6.0, 0.4,
    size=16, bold=True, color=ORANGE)
txt(s, "Pipeline Comparison Study", 6.9, 1.22, 6.0, 0.3,
    size=12, color=LGRAY, italic=True)
bullet_block(s, [
    "THREE interchangeable pipelines (toggle flags)",
    "Sinusoidal aberration injected as ground truth",
    "Monte Carlo sweep over defocus dz × spatial freq v0",
    "",
    "DEFOCUS_MODEL flag:",
    "  pupil_zernike | image_AS | my_fdpr",
    "",
    "PUPIL_READOUT flag:",
    "  mft_rev | hcipy_backward | fft",
    "",
    "Noise optional (if_noise flag)",
    "Convergence snapshots saved during iteration",
], left=6.9, top=1.55, width=6.0, height=4.8, size=13, bullet_color=LGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Pipeline 1: pupil_zernike + mft_rev
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "Pipeline 1 — pupil_zernike + mft_rev",
             "Classic Zernike defocus in the pupil plane · FocusDiversePhaseRetrieval solver")

# forward model
box(s, 0.3, 0.75, 8.0, 2.5, fill_color=RGBColor(0x08,0x18,0x30), line_color=BLUE, line_width=Pt(1))
txt(s, "Forward Model (Defocus Simulation)", 0.4, 0.8, 7.8, 0.35,
    size=15, bold=True, color=BLUE)
code_block(s, [
    "phi_def = calculate_defocus_phase(seal_parameters, dz_mm, ...)",
    "# Zernike mode 4 (focus), scaled to desired P2V in radians",
    "",
    "wf_focus = Wavefront(pupil * exp(i·φ_sine), λ)",
    "wf_defoc = Wavefront(pupil * exp(i·(φ_sine + φ_def)), λ)",
    "",
    "psf_focus = prop_p2f(wf_focus).power.shaped    # Fraunhofer",
    "psf_defoc = prop_p2f(wf_defoc).power.shaped",
], left=0.35, top=1.18, width=7.9, height=2.0)

# phase retrieval
box(s, 0.3, 3.4, 8.0, 2.5, fill_color=RGBColor(0x08,0x18,0x30), line_color=CYAN, line_width=Pt(1))
txt(s, "Phase Retrieval — FocusDiversePhaseRetrieval (angular spectrum)", 0.4, 3.45, 7.8, 0.35,
    size=15, bold=True, color=CYAN)
code_block(s, [
    "mp = FocusDiversePhaseRetrieval([psf0, psfd], λ_um, [dx_um], [dz_um])",
    "# image_sharpening library: GS loop using angular-spectrum propagation",
    "# internally matches defocus by free-space AS transfer function",
    "",
    "for it in range(250):",
    "    psf_rec = mp.step()   # returns complex focal-plane field",
    "",
    "raw_pupil = angle(mft_rev(psf_rec, conf))    # Matrix FT back to pupil",
    "real_pupil = resize(raw_pupil, (256,256)) * telescope_pupil.shaped",
], left=0.35, top=3.83, width=7.9, height=2.0)

# notes column
bullet_block(s, [
    "Zernike defocus is analytically clean",
    "Consistent with HCIPy forward model",
    "mft_rev uses geometry from InstrumentConfiguration",
    "Most physically transparent pipeline",
    "Used in FDPRno_phase_error.py as sole pipeline",
], left=8.5, top=0.75, width=4.5, height=5.8,
   title="Notes", fill=RGBColor(0x08,0x18,0x30), border=BLUE, size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Pipeline 2: image_AS + hcipy_backward
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "Pipeline 2 — image_AS + hcipy_backward",
             "Angular spectrum propagation in the image plane · HCIPy adjoint readout")

box(s, 0.3, 0.75, 8.0, 2.8, fill_color=RGBColor(0x08,0x18,0x30), line_color=ORANGE, line_width=Pt(1))
txt(s, "Forward Model — Angular Spectrum in Image Plane", 0.4, 0.8, 7.8, 0.35,
    size=15, bold=True, color=ORANGE)
code_block(s, [
    "# Step 1: focused PSF via HCIPy Fraunhofer (same as Pipeline 1)",
    "wf_focus_focal = prop_p2f(wf_focus)           # complex focal field",
    "E_focus = wf_focus_focal.electric_field.shaped",
    "",
    "# Step 2: propagate focal E-field through free space by +dz",
    "H_fwd = _angular_spectrum_transfer_function(E.shape, λ_um, dx_um, dz_um)",
    "E_def  = _angular_spectrum_prop(E_focus, H_fwd)",
    "psf_defoc = |E_def|²",
    "",
    "# Power is renormalized to match focused PSF total photons",
], left=0.35, top=1.18, width=7.9, height=2.35)

box(s, 0.3, 3.7, 8.0, 2.5, fill_color=RGBColor(0x08,0x18,0x30), line_color=GREEN, line_width=Pt(1))
txt(s, "Pupil Readout — HCIPy Adjoint Fraunhofer", 0.4, 3.75, 7.8, 0.35,
    size=15, bold=True, color=GREEN)
code_block(s, [
    "# prop_p2f.backward() = adjoint of the forward propagator",
    "# Avoids the geometry assumptions baked into mft_rev",
    "",
    "wf_focal = Wavefront(Field(E_rec.ravel(), focal_grid), λ)",
    "wf_pupil = prop_p2f.backward(wf_focal)",
    "real_pupil = angle(wf_pupil.electric_field.shaped) * pupil_mask",
    "",
    "# Note: piston offset from normalisation factor must be corrected separately",
], left=0.35, top=4.13, width=7.9, height=2.0)

bullet_block(s, [
    "Models defocus as physical free-space propagation in the image plane",
    "More realistic for lab use (camera physically moved)",
    "AS internally called by image_sharpening — conflicts with mft_rev → must use hcipy_backward",
    "Normalisation piston artefact needs correction",
    "Most complex pipeline to interpret",
], left=8.5, top=0.75, width=4.5, height=5.8,
   title="Notes", fill=RGBColor(0x08,0x18,0x30), border=ORANGE, size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Pipeline 3: my_fdpr (SimpleFDPR) + fft
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "Pipeline 3 — my_fdpr (SimpleFDPR) + FFT",
             "Custom Gerchberg-Saxton with FFT — no external angular-spectrum dependency")

box(s, 0.3, 0.75, 8.0, 1.4, fill_color=RGBColor(0x08,0x18,0x30), line_color=GREEN, line_width=Pt(1))
txt(s, "Forward Model — same PSFs as Pipeline 1 (Zernike defocus)", 0.4, 0.8, 7.8, 0.35,
    size=15, bold=True, color=GREEN)
code_block(s, [
    "psf_focus, psf_defoc  ←  identical to pupil_zernike  (Fraunhofer + Zernike)",
], left=0.35, top=1.18, width=7.9, height=0.85)

box(s, 0.3, 2.3, 8.0, 3.95, fill_color=RGBColor(0x08,0x18,0x30), line_color=CYAN, line_width=Pt(1))
txt(s, "Phase Retrieval — SimpleFDPR (custom Gerchberg-Saxton, FFT-only)", 0.4, 2.35, 7.8, 0.35,
    size=15, bold=True, color=CYAN)
code_block(s, [
    "mp = SimpleFDPR(psf0)",
    "phi_def_resized = resize(phi_def.shaped, psf0.shape)",
    "mp.add_defocused_image(psfd, phi_def_resized)",
    "",
    "# Inside mp.step():",
    "focal_field = A_focus * exp(i·φ_estimate)",
    "pupil_field  = ifft2(focal_field)                   # pupil via FFT",
    "pupil_defoc  = |pupil| * exp(i·(∠pupil + φ_def))   # add defocus",
    "focal_defoc  = fft2(pupil_defoc)                    # to focal",
    "focal_defoc  = A_defoc * exp(i·∠focal_defoc)        # enforce amplitude",
    "pupil_back   = ifft2(focal_defoc)                   # back to pupil",
    "pupil_focus  = |pupil_back| * exp(i·(∠pupil_back - φ_def))  # remove defocus",
    "focal_field  = fft2(pupil_focus)                    # to focal",
    "focal_field  = A_focus * exp(i·∠focal_field)        # enforce focused amplitude",
    "",
    "# Readout: direct ifft2 of reconstructed focal field",
    "pupil_phase = mp.get_pupil_phase()   # = ∠(ifft2(A_focus·e^iφ))",
], left=0.35, top=2.73, width=7.9, height=3.45)

bullet_block(s, [
    "Fully self-contained — no image_sharpening dependency",
    "FFT is less accurate than MFT for non-square/non-Nyquist grids",
    "Defocus applied as phase in pupil (same physical model as P1)",
    "Readout via direct FFT — geometry consistent with solver",
    "Easier to modify/debug",
    "250 iterations, cost tracked per defocused image",
], left=8.5, top=0.75, width=4.5, height=5.8,
   title="Notes", fill=RGBColor(0x08,0x18,0x30), border=GREEN, size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Solver Deep Dive: FocusDiversePhaseRetrieval vs SimpleFDPR
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "Phase Retrieval Solvers: Head-to-Head",
             "FocusDiversePhaseRetrieval (image_sharpening)  vs  SimpleFDPR (custom FFT)")

# left
box(s, 0.3, 0.75, 6.0, 6.3, fill_color=RGBColor(0x08,0x18,0x30), line_color=CYAN, line_width=Pt(2))
txt(s, "FocusDiversePhaseRetrieval", 0.4, 0.82, 5.8, 0.38,
    size=16, bold=True, color=CYAN)
bullet_block(s, [
    "From image_sharpening library",
    "Angular-spectrum propagation between focus positions",
    "Propagates in the IMAGE plane (physical dz in µm)",
    "Input: list of N+1 PSFs, N pixel scales, N defocus distances",
    "Each step() returns complex focal-plane E-field",
    "Cost function per PSF tracked in mp.cost_functions",
    "mft_rev readout (Matrix Fourier Transform) — more accurate sampling",
    "Used in: FDPRno_phase_error.py and Pipelines 1 & 2",
], left=0.4, top=1.25, width=5.8, height=5.5, size=13, bullet_color=LGRAY)

# right
box(s, 7.0, 0.75, 6.0, 6.3, fill_color=RGBColor(0x08,0x18,0x30), line_color=GREEN, line_width=Pt(2))
txt(s, "SimpleFDPR", 7.1, 0.82, 5.8, 0.38,
    size=16, bold=True, color=GREEN)
bullet_block(s, [
    "Defined directly in fdpr_vs_AS_save.py",
    "Defocus applied as a PUPIL-PLANE phase offset",
    "Propagation uses FFT (fft2 / ifft2) throughout",
    "Input: focused PSF + list of (psf_defoc, defocus_phase) pairs",
    "Each step() returns real-valued PSF intensity",
    "Cost tracked per defocused image in data['cost']",
    "FFT readout (get_pupil_phase) — consistent with solver",
    "Used in: Pipeline 3 (my_fdpr) only",
], left=7.1, top=1.25, width=5.8, height=5.5, size=13, bullet_color=LGRAY)

# divider
box(s, 6.55, 0.75, 0.08, 6.3, fill_color=BLUE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Pupil Readout Methods
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "Pupil Readout Methods",
             "After FDPR converges: how do we get from focal-plane E-field to pupil phase?")

# three cards
for idx, (title, color, lines, note) in enumerate([
    ("mft_rev", CYAN,
     ["raw = ∠( mft_rev(psf_rec, conf) )",
      "pupil = resize(raw, 256×256) × pupil_mask",
      "",
      "Matrix Fourier Transform — samples any",
      "region of the pupil at arbitrary resolution.",
      "Geometry encoded in InstrumentConfiguration."],
     "Paired with: pupil_zernike\nMost accurate for non-uniform grids"),
    ("hcipy_backward", GREEN,
     ["wf_focal = Wavefront(E_rec, focal_grid, λ)",
      "wf_pupil = prop_p2f.backward(wf_focal)",
      "pupil = ∠(wf_pupil.electric_field.shaped)",
      "",
      "HCIPy adjoint Fraunhofer propagator.",
      "Physically consistent with forward model.",
      "Piston term requires correction."],
     "Paired with: image_AS\nAvoids mft_rev geometry conflict"),
    ("fft", ORANGE,
     ["pupil = mp.get_pupil_phase()",
      "  = ∠( ifft2( A_focus · e^iφ_est ) )",
      "",
      "Direct inverse FFT of estimated focal field.",
      "Fast, simple, no extra dependencies.",
      "Geometry must match FFT assumptions."],
     "Paired with: my_fdpr\nConsistent with SimpleFDPR solver"),
]):
    left = 0.3 + idx * 4.35
    box(s, left, 0.75, 4.1, 5.8,
        fill_color=RGBColor(0x08,0x18,0x30), line_color=color, line_width=Pt(2))
    txt(s, title, left+0.1, 0.82, 3.9, 0.38,
        size=17, bold=True, color=color)
    code_block(s, lines, left+0.05, 1.25, 4.0, 3.2)
    txt(s, note, left+0.1, 4.55, 3.9, 1.85,
        size=12, color=LGRAY, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Pipeline Combinations & Compatibility
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "Pipeline Combinations & Compatibility",
             "Why certain pairings are required — and what breaks when mixed")

# table header
for col, label, color in [
    (0.3,  "DEFOCUS_MODEL",   CYAN),
    (3.8,  "Solver",          CYAN),
    (6.5,  "PUPIL_READOUT",   CYAN),
    (9.2,  "Why paired",      CYAN),
]:
    box(s, col, 0.9, 3.2 if col < 9 else 4.0, 0.45, fill_color=BLUE)
    txt(s, label, col+0.08, 0.95, 3.1, 0.35,
        size=14, bold=True, color=WHITE)

rows = [
    ("pupil_zernike",  "FocusDiversePhaseRetrieval\n(image_sharpening)",  "mft_rev",
     "Zernike defocus and mft_rev share the same HCIPy/MFT geometry. Internally consistent."),
    ("image_AS",       "FocusDiversePhaseRetrieval\n(image_sharpening)",  "hcipy_backward",
     "AS is used inside image_sharpening; mft_rev geometry conflicts. HCIPy adjoint matches forward."),
    ("my_fdpr",        "SimpleFDPR\n(custom FFT GS)",                     "fft",
     "SimpleFDPR uses FFT throughout. Direct FFT readout is the only consistent choice."),
]

for r, (dm, solver, ro, why) in enumerate(rows):
    top = 1.5 + r * 1.6
    row_color = [RGBColor(0x0C,0x20,0x3A), RGBColor(0x0A,0x1A,0x30), RGBColor(0x0C,0x20,0x3A)][r]
    accent = [CYAN, GREEN, ORANGE][r]
    for col, val, w in [(0.3, dm, 3.3), (3.8, solver, 2.5), (6.5, ro, 2.5), (9.2, why, 4.0)]:
        box(s, col, top, w, 1.5, fill_color=row_color, line_color=accent, line_width=Pt(1))
        txt(s, val, col+0.1, top+0.1, w-0.2, 1.3, size=12, color=LGRAY)

# warning note
box(s, 0.3, 6.2, 12.7, 0.7, fill_color=RGBColor(0x30,0x10,0x00), line_color=ORANGE, line_width=Pt(1))
txt(s, "⚠  Mixing pipelines (e.g. image_AS + mft_rev) will silently produce wrong results — "
    "the angular spectrum call inside image_sharpening conflicts with the MFT geometry assumptions in mft_rev.",
    0.4, 6.28, 12.5, 0.55, size=12, color=ORANGE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Monte Carlo Setup (fdpr_vs_AS only)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "Monte Carlo Setup — fdpr_vs_AS_save.py",
             "Systematic sweep over defocus distance and sinusoidal spatial frequency")

bullet_block(s, [
    "Outer loop:  dz  (defocus distance, mm)  — linspace(5, 250, N)",
    "Inner loop:  v0  (spatial frequency, cycles/aperture)  — linspace(0.5, 30, N)",
    "N_trials random noise realisations per (dz, v0) grid point",
    "",
    "At each point, N_trials independent PSF pairs are generated and run through FDPR",
    "RMS residual  =  truth sinusoid − median-subtracted reconstruction  [nm]",
    "RMS total     =  magnitude of reconstruction  [nm]",
    "Convergence rate  =  fraction of trials that produce a finite result",
    "",
    "Optional noise:  Gaussian read noise  σ_e = 11 e⁻/px  (if_noise flag)",
    "Snapshot images saved at iterations: 0, 10, 20, 50, 100, 150, 200, 250, …",
    "Results saved to .npz: residual/total RMS arrays + convergence rate grid",
], left=0.4, top=0.85, width=7.8, height=5.9,
   title="Grid & Metrics", fill=RGBColor(0x08,0x18,0x30), border=BLUE, size=13)

bullet_block(s, [
    "verbose = True",
    "  → 2×4 plot at first & last trial of every grid point",
    "  → Shows: injected phase, PSFs, reconstruction, difference",
    "",
    "do_spot_check = True",
    "  → Full 4×4 diagnostic plot at a single chosen (dz, v0)",
    "  → Includes error histogram & central slice",
    "",
    "snapshot = True",
    "  → PSF, pupil phase, & cost plotted at each saved iteration",
    "  → Useful for diagnosing convergence rate",
], left=8.4, top=0.85, width=4.7, height=5.9,
   title="Diagnostic Flags", fill=RGBColor(0x08,0x18,0x30), border=CYAN, size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FDPRno_phase_error Pipeline Detail
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "FDPRno_phase_error.py — Pipeline Detail",
             "Single-run baseline: what does FDPR reconstruct when there is nothing to find?")

box(s, 0.3, 0.75, 12.7, 2.65, fill_color=RGBColor(0x08,0x18,0x30), line_color=CYAN, line_width=Pt(1))
txt(s, "Step-by-Step", 0.4, 0.82, 12.5, 0.35, size=15, bold=True, color=CYAN)
code_block(s, [
    "# 1. Flat (aberration-free) pupil",
    "focused_wf = Wavefront(pupil * exp(i·0), λ)        # zero phase",
    "",
    "# 2. Seven defocus positions: 0, ±4, ±8, ±12 mm",
    "for dz_mm in [0, -12, -8, -4, +4, +8, +12]:",
    "    phi_def = calculate_defocus_phase(params, dz_mm)   # Zernike",
    "    psf = prop_p2f( Wavefront(pupil * exp(i·phi_def), λ) ).intensity",
    "    psf_list.append(psf)",
    "",
    "# 3. FDPR: 200 iterations with FocusDiversePhaseRetrieval",
    "mp = FocusDiversePhaseRetrieval(psf_list, λ_um, dx_list, defocus_um_list)",
    "for _ in range(200): psf_rec = mp.step()",
    "",
    "# 4. Readout via mft_rev, resize, mask",
    "raw = angle(mft_rev(psf_rec, conf))",
    "pupil_phase = resize(raw, (256,256)) * telescope_pupil.shaped",
], left=0.35, top=1.15, width=12.6, height=2.2)

bullet_block(s, [
    "Truth = zero phase  →  any reconstruction IS the error",
    "Median subtraction removes piston: med_subtracted = pupil_phase − median(phase inside pupil)",
    "Residual = focused_wavefront_pupil.phase − med_subtracted  ≈  0 − med_subtracted",
    "RMS and P2V reported in both radians and nm  (× λ/2π × 10⁹)",
    "Cost functions plotted per defocus position (6 curves on semilogy plot)",
], left=0.4, top=3.5, width=12.5, height=2.0,
   title="Error Analysis", fill=RGBColor(0x08,0x18,0x30), border=GREEN, size=13)

bullet_block(s, [
    "1D grid section (oned_grid=True): single dz=5mm, sinusoidal aberration injected",
    "Sweeps spatial frequency v0, plots RMS vs v0 — quick sensitivity check",
    "Provides a minimal version of the full MC sweep in fdpr_vs_AS_save.py",
], left=0.4, top=5.6, width=12.5, height=1.6,
   title="Bonus: 1D Grid Section", fill=RGBColor(0x08,0x18,0x30), border=ORANGE, size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary Comparison Table
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
slide_header(s, "Summary Comparison", "All three pipelines + baseline at a glance")

# table
headers = ["Aspect", "FDPRno_phase_error", "P1: pupil_zernike", "P2: image_AS", "P3: my_fdpr"]
col_w   = [2.6, 2.5, 2.1, 2.1, 2.1]
col_x   = [0.3]
for w in col_w[:-1]: col_x.append(col_x[-1] + w + 0.06)

rows_data = [
    ("Injected aberration",   "None (flat phase)",      "Sinusoid (φ_sine)", "Sinusoid (φ_sine)", "Sinusoid (φ_sine)"),
    ("Defocus method",        "Zernike pupil-plane",    "Zernike pupil-plane","Angular spectrum\n(image plane)","Zernike pupil-plane"),
    ("Defocus positions",     "6  (±4, ±8, ±12 mm)",   "1 per MC point",    "1 per MC point",    "1 per MC point"),
    ("FDPR solver",           "FocusDiversePhaseRetr.", "FocusDiversePhaseRetr.","FocusDiversePhaseRetr.","SimpleFDPR (FFT)"),
    ("Pupil readout",         "mft_rev",                "mft_rev",           "hcipy_backward",    "fft"),
    ("Monte Carlo",           "No",                     "Yes (dz × v0 grid)","Yes (dz × v0 grid)","Yes (dz × v0 grid)"),
    ("Noise modelled",        "No",                     "Optional",          "Optional",          "Optional"),
    ("Primary purpose",       "Noise floor baseline",   "Pipeline comparison","Pipeline comparison","Dependency-free test"),
]

row_h = 0.56
for ci, (header, cw, cx) in enumerate(zip(headers, col_w, col_x)):
    box(s, cx, 0.82, cw, 0.42, fill_color=BLUE)
    txt(s, header, cx+0.06, 0.88, cw-0.12, 0.3,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows_data):
    top = 1.28 + ri * row_h
    bg_c = RGBColor(0x0C,0x20,0x3A) if ri % 2 == 0 else RGBColor(0x08,0x16,0x2C)
    for ci, (val, cw, cx) in enumerate(zip(row, col_w, col_x)):
        box(s, cx, top, cw, row_h-0.04,
            fill_color=bg_c, line_color=RGBColor(0x20,0x40,0x60), line_width=Pt(0.5))
        txt(s, val, cx+0.06, top+0.04, cw-0.12, row_h-0.1,
            size=10, color=LGRAY if ci > 0 else WHITE, bold=(ci==0))

# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
out = "/Users/joshuapotter/Documents/SEAL/FDPRNotebooks/FDPR_Pipeline_Comparison.pptx"
prs.save(out)
print(f"Saved: {out}")
